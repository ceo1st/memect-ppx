import io
import logging



import statistics
from typing import Final, Sequence


from pptx import Presentation, presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR,PP_ALIGN
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Pt

import pptx_ea_font


from .font import FontInfo, FontParser

from memect.base.bbox import BBox



from .base import KColor, KDocument, KFigure, KFormula, KLine, KObject, KPage, KRect, KTable, KText, KTextbox, KTextline


class PptxBuilder:
    """
    如果需要嵌入字体，只能够嵌入可以免费商用的，如：
    思源黑体，阿里普惠体，鸿蒙黑体。

    如果是在window上查看，通常系统自带了微软雅黑

    所以：统一使用思源黑体，支持linux/mac/windows，用户的电脑上面没有，需要自行安装。

    """
    _logger=logging.getLogger(f'{__module__}.{__qualname__}')
    def __init__(self):
        super().__init__()
        self._font_parser = FontParser()
    
    def build(self,doc:KDocument)->bytes:
        assert len(doc.pages)>0
        prs = Presentation()
        # pptx的slide必须使用同一的size，不能够使用不同的
        # 典型的就是A4纸，这里可以使用Inches，或者
        # Pt(566) Pt(842) = 8.27*72  11.69*72

        #正常的文档页面大小都是一样的，可能有些例外，如：横版了
        #如：两个页面，500*800，800*500 => 800*800
        #
        width = max(p.width for p in doc.pages)
        height = max(p.height for p in doc.pages)
        prs.slide_width = Pt(width)
        prs.slide_height = Pt(height)
        for page in doc.pages:
            if page.skipped:
                #TODO 创建一个空白页？
                pass
            else:
                self._render_page(prs, page)
        
        with io.BytesIO() as fp:
            prs.save(fp)
            return fp.getvalue()

    def _render_page(self, prs: presentation.Presentation, page: KPage):
        # 空白页面
        # 使用Pt单位，刚好和pdf的一致，都是1inch=72pt，就不需要使用inch单位了
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        #可以先渲染了其他，获得对应的背景颜色？
        #因为表格使用了自定义的style，所以就需要把表格区域的line，rect去掉

        
        #self._render_block(slide,adjust_other())
        
        if page.header.objects:
            self._render_block(slide,page.header.objects)
        
        #if page.footnotes:
            #for footnote in page.footnotes:
                #self._render_block(slide,footnote.body)
        if page.footer.objects:
            self._render_block(slide,page.footer.objects)
        
        self._render_block(slide,page.objects)

    def _render_block(self,slide:Slide,objects:Sequence[KObject]):
        for obj in objects:
            if not obj.bbox.is_valid():
                #如：deepseek返回的内容，存在无法获得bbox的情况，为了保留完整的markdown内容，仍然创建对象
                #但是无法在pptx中渲染，因为不知道坐标
                self._logger.warning('第%s页，对象的bbox无效，不显示该对象,object=%s,bbox=%s',obj,obj.bbox)
                continue
            if isinstance(obj,KText):
                #文本，没有每个字符的坐标
                self._render_text(slide,obj)
            elif isinstance(obj,KTextbox):
                pass
            elif isinstance(obj, KFigure):
                self._render_figure(slide,obj)
            elif isinstance(obj, KFormula):
                self._render_formula(slide,obj)
            elif isinstance(obj, KTable):
                self._render_table(slide,obj)
            elif isinstance(obj,KRect):
                self._render_rect(slide,obj)
            elif isinstance(obj,KLine):
                self._render_line(slide,obj)
            else:
                pass

    def _render_rect(self,slide:Slide,rect:KRect):
        m = rect.page.to_lt()
        x0,y0,x1,y1 = rect.bbox.transform(m)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Pt(x0),Pt(y0),
            Pt(x1-x0),Pt(y1-y0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._to_color(rect.color)
        shape.line.fill.background()
    
    def _render_line(self,slide:Slide,line:KLine):
        m = line.page.to_lt()
        #TODO 这里使用bbox可能不合适了
        x0,y0,x1,y1 = line.bbox.transform(m)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.LINE_INVERSE,
            Pt(x0),Pt(y0),
            Pt(x1-x0),Pt(y1-y0)
        )
        shape.line.color.rgb = self._to_color(line.color)
        shape.line.width = Pt(line.width)

    def _render_text(self,slide:Slide,text:KText,*,text_frame:TextFrame|None=None):
        
        def get_fontsize(bbox:BBox,text:str,ratio:float=1.5)->float:
            """
            计算合适的字体大小，使文本刚好填充 bbox

            Args:
                bbox: 文本区域
                text: 文本内容
                ratio: 行距倍数

            Returns:
                字体大小（pt）
            """
            if not text:
                return 12.0

            # PowerPoint 行高公式: line_height = (1.2018 × ratio + 0.0034) × font_size
            line_height_factor = 1.2018 * ratio + 0.0034

            # 根据实际文本计算平均字符宽度系数
            # 中文/全角字符 ≈ 1.0 × font_size，英文/半角字符 ≈ 0.5 × font_size
            total_width_factor = 0.0
            for ch in text:
                if ord(ch) > 0x7F:
                    total_width_factor += 1.0
                else:
                    total_width_factor += 0.5
            avg_char_width_factor = total_width_factor / len(text)

            # 二分查找合适的字体大小
            min_size, max_size = 1.0, 200.0
            best_size = 12.0

            for _ in range(20):  # 最多迭代20次
                font_size = (min_size + max_size) / 2

                # 计算每行可容纳的字符数
                avg_char_width = font_size * avg_char_width_factor
                chars_per_line = max(1, int(bbox.width / avg_char_width))

                # 估算行数（简单按字符数除以每行字符数）
                estimated_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)

                # 计算所需总高度
                line_height = line_height_factor * font_size
                required_height = line_height * estimated_lines

                # 判断是否合适
                if abs(required_height - bbox.height) < 1.0:  # 误差小于1pt
                    best_size = font_size
                    break
                elif required_height > bbox.height:
                    # 字体太大，减小
                    max_size = font_size
                else:
                    # 字体太小，增大
                    min_size = font_size
                    best_size = font_size

            return best_size

        s = text.text
        m:Final=text.page.to_lt()
        bbox:Final= text.bbox.transform(m)
        if text_frame is None:
            x0,y0,x1,y1=bbox
            textbox = slide.shapes.add_textbox(
                Pt(x0), Pt(y0),
                Pt(x1-x0), Pt(y1-y0)
            )
            text_frame = textbox.text_frame
            text_frame.clear()
        
        text_frame.margin_bottom=Pt(0)
        text_frame.margin_top=Pt(0)
        text_frame.margin_left=Pt(0)
        text_frame.margin_right=Pt(0)
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.word_wrap = True
        
        #行距的倍数
        line_spacing=1.2
        text_image = text.page.crop(text.bbox)
        assert text_image is not None
        font_info = self._font_parser.parse(text_image,s,full_image_size=text.page.image.size)
        #因为没有每个字符的坐标，所以就做一个简单的假设，按行书写，然后自动换行，高度刚好充满，
        #所以，只需要计算出一个合理的字体大小，就大体可以还原
        #没有要求直接传递_Paragraph是因为这个类型是私有的
        if not text_frame.paragraphs:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[-1]
        p.space_after=Pt(0)
        p.space_before=Pt(0)
        p.line_spacing=line_spacing
        #p.alignment=
        run = p.add_run()
        run.font.size=Pt(get_fontsize(bbox,s,ratio=line_spacing))
        run.font.bold = font_info.bold
        run.font.color.rgb = RGBColor(*font_info.color)
        run.font.name=font_info.family
        run.text = s
        

    def _render_textbox(self, slide: Slide, text:KTextbox,*,text_frame:TextFrame|None=None):
        # Textbox with formatted text

        #为了简化，对于Text，仅仅识别一个字体即可，即使多行，因为对于复杂的情况，可能每一个字符都使用不同的字体
        #获得文字最长的一行进行识别即可
        #暂时没有在image2json中去获得字体的颜色，信息等，因为不需要
        def get_text_font(text:Text)->FontInfo:
            #这里是从视觉上获得信息，而不考虑是否为pdf？
            if text.page.type==PageType.PDF:
                #可以根据字体信息来判断是否为serif/bold，只是
                #line=sorted(text.lines,key=lambda line:len(line.text2),reverse=True)[0]
                #这里可以随便了，因为只是用来构造
                #return self._font_parser._get_font(serif=serif,bold=bold)
                pass

            #获得最长一行的
            line=sorted(text.lines,key=lambda line:len(line.text2),reverse=True)[0]
            return self._font_parser.parse(text.page.crop(line.bbox),line.text,full_image_size=text.page.pil_image.size)

        
        if text_frame is None:
            #如果没有指定，那么支持旋转字符
            texts:list[Text]=[]
            i=0
            for j,line in enumerate(text.lines):
                if len(line.objects)==1 and len(line.chars)==1 and line.chars[0].rotate!=0:
                    #表示有旋转字符
                    if j>i:
                        t  = text.slice(i,j)
                        assert t is not None
                        texts.append(t)
                    t = text.slice(j,j+1)
                    assert t is not None
                    texts.append(t)
                    i=j+1
                else:
                    pass
            
            if i<len(text.lines):
                t = text.slice(i)
                assert t is not None
                texts.append(t)
            else:
                pass
            
            for t in texts:
                
                #TODO 使用具体的字体渲染的后的真实宽度可能溢出
                #因为pdf使用的字体可能不一致，即使字体大小一致（高度），宽度也可能不一致
                #TODO 如果包含有旋转的字符，就必须使用新的textbox，也就是需要创建多个

                rotate = 0
                if len(t.objects)==1 and len(t.chars)==1 and t.chars[0].rotate!=0:
                    rotate = t.chars[0].rotate
                    #获得没有旋转之前的坐标
                    bbox = self._get_unrotated_bbox(t.chars[0])
                else:
                    bbox = t.bbox
                x0, y0, x1, y1 = bbox.change_origin(t.page.height)
                textbox = slide.shapes.add_textbox(
                    Pt(x0), Pt(y0),
                    Pt(x1-x0), Pt(y1-y0)
                )
                if rotate!=0:
                    textbox.rotation=rotate
                text_frame = textbox.text_frame
                text_frame.clear()
                self._render_text(slide,t,text_frame=text_frame)
        else:
            text_frame.margin_bottom=Pt(0)
            text_frame.margin_top=Pt(0)
            text_frame.margin_left=Pt(0)
            text_frame.margin_right=Pt(0)
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            text_frame.word_wrap = False
            

            #如果需要支持编辑修改，使用下面这个更好，因为改变输入的时候，自动改变shape
            #默认为MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT（也就是自动调整shape适配text的宽度）
            #text_frame.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            #text_frame.auto_size=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            # 和docx类似，都是<p><run>xx</run></p>
            # textline => <p>
            # span => <run>

            #有两种实现，整个文本只使用一个<p>，然后通过line_break换行，通过line_spaceing来控制行间距
            #第二种就是一行就是一个<p>，通过space_after控制间距

            #在同一个Text中，目前使用相同的字体
            #计算出fontsize，line height（行间距，通常为1，1.2，1.5等）
            font = get_text_font(text)
            fontsizes:list[int]=[]
            height_ratios:list[float]=[]
            for i,line in enumerate(text.lines):
                fontsize = font.get_fontsize(line.text,line.bbox)
                fontsizes.append(fontsize)
                fb=font.get_fontbbox(line.text,fontsize)
                #print('==>fff',fontsize,line.text,line.bbox,fb)
                if i>0:
                    line_height=text.lines[i-1].bbox.y0-line.bbox.y0
                    #可能有点重叠
                    line_height=max(line_height,line.bbox.height)
                    height_ratios.append(line_height/max(line.bbox.height,fb.height))

            #可以使用平均值或者中位数，如果两行之间高度差别大的，就不使用
            #20-16 
            
            sorted_fontsizes=sorted(fontsizes)
            #差别过大会导致width过大，溢出
            #使用最小的，会导致宽度太小，有空余
            #使用最大的，会导致宽度太大，溢出
            #使用中位数，如果大的还是差别过大
            #先获得中位数，然后计算是否宽度溢出

            if (sorted_fontsizes[-1]-sorted_fontsizes[0])/sorted_fontsizes[-1]<=0.2:
                mid_fontsize=statistics.median_high(fontsizes)
                diff_count=0
                for line,fontsize in zip(text.lines,fontsizes):
                    if abs(fontsize-mid_fontsize)<=2:
                        continue
                    fontbbox = font.get_fontbbox(line.text,mid_fontsize)
                    if fontbbox.width-line.bbox.width>=fontsize:
                        #如果超过了一个字符的宽度，就使用原始的fontsize
                        diff_count+=1
                    else:
                        pass
                
                if diff_count==0:
                    #如果都在范围内，使用相同的
                    fontsizes=[mid_fontsize]*len(fontsizes)
                else:
                    #如果有溢出的，就使用各自的fontsize
                    pass
            else:
                #差别过大，不调整了
                pass

            #计算行间距
            height_ratio=1
            #print('==========>>',text.text,height_ratios)
            if len(height_ratios)>0:
                height_ratio = statistics.mean(height_ratios)
                if height_ratio<=1.1:
                    height_ratio=1
                elif height_ratio<=1.4:
                    height_ratio=1.2
                elif height_ratio<=1.8:
                    height_ratio=1.5
                elif height_ratio<=2.4:
                    height_ratio=2
                else:
                    height_ratio=2.5
                    #height_ratio = round(height_ratio,1)
            

            #或者根据现在的fontsize
            #(text.bbox.height-sum(fontsizes))/len(text.lines)

            align=PP_ALIGN.LEFT
            if len(text.lines)>1:
                #计算左对齐，右对齐，居中支持
                lines = sorted(text.lines,key=lambda line:line.bbox.width)
                min_line = lines[0]
                max_line = lines[-1]
                #现在仅仅考虑是否居中
                #第一行可能突出/缩进
                #------max------
                #  ----min----
                d1=min_line.bbox.x0-max_line.bbox.x0
                d2=max_line.bbox.x1-min_line.bbox.x1
                if abs(d1-d2)<=10 and min(d1,d2)>10:
                    align=PP_ALIGN.CENTER
                pass

            for i, line in enumerate(text.lines):
                if i>=len(text_frame.paragraphs):
                    text_frame.add_paragraph()
                if i>0:
                    text_frame.paragraphs[-1].line_spacing=height_ratio
                text_frame.paragraphs[-1].alignment=align
                text_frame.paragraphs[-1].space_before=Pt(0)
                #print('rrrr',line.text,fontsizes[i])
                self._render_textline(slide,text_frame,i,line,font,fontsizes[i])


    def _render_textline(self, slide: Slide,text_frame:TextFrame,index:int, line: KTextline,font:FontInfo,fontsize:int):
        page:Final = line.page
        class _Span:
            def __init__(self):
                super().__init__()
                self.chars: list[Char] = []

            @property
            def text(self) -> str:
                return ''.join(c.text for c in self.chars)
            
            @property
            def bbox(self)->BBox:
                span_bbox = BBox.x_union(self.chars)
                assert span_bbox is not None 
                return span_bbox

            def accept(self, c: Char) -> bool:
                c1 = self.chars[0]
                if c1.alike(c):
                    return True
                else:
                    return False

        def split_objects(objects: Sequence[TObject]) -> list[_Span | TObject]:
            span: _Span | None = None
            groups: list[_Span | TObject] = []
            for obj in objects:
                if isinstance(obj, Char):
                    if span is None or not span.accept(obj):
                        span = _Span()
                        groups.append(span)
                    span.chars.append(obj)
                else:
                    span = None
                    groups.append(obj)
            return groups


        #没有要求直接传递_Paragraph是因为这个类型是私有的
        if index>=len(text_frame.paragraphs):
            text_frame.add_paragraph()
        p = text_frame.paragraphs[index]
        
        for obj in split_objects(line.objects):
            # TODO 如果需要考虑斜体粗体等支持，就可以先划分为[figure,formul,span1,span2]
            # if isinstance(obj,Char):
            
            if isinstance(obj, _Span):
                span = obj
                char = span.chars[0]
                s = span.text
                run = p.add_run()
                run.text = s

                #TODO 如果旋转了，bbox是旋转的bbox，需要先获得
                #fontsize = get_fontsize(span)
                run.font.size=Pt(fontsize)
                
                if page.type==PageType.PDF:
                    #如果是来自pdf的，使用pdf解析处理的粗体/斜体/下划线？
                    run.font.bold = char.bold
                    run.font.underline = char.underline
                    run.font.italic = char.italic
                    #颜色不使用pdf解析出来的，因为可能使用了背景
                    run.font.color.rgb=self._to_color(char.color)
                    #run.font.color.rgb=RGBColor(*font.color)
                    #还是使用ocr识别的字体而不是来自pdf的
                    family = font.family
                else:
                    run.font.bold = font.bold
                    run.font.color.rgb = RGBColor(*font.color)
                    family=font.family
                run.font.name=family
                #设置这个才会自动设置中文的字体，否则默认仅仅设置latin的字体，使用powerpoint打开ppt没有问题，使用wps有
                pptx_ea_font.set_font(run,font.family)
                if char.strickout:
                    #目前没有api支持
                    #run.font._element.attrib['strike']='sngStrike'
                    pass

            elif isinstance(obj, Figure):
                self._render_figure(slide,obj)
            elif isinstance(obj, Formula):
                self._render_formula(slide,obj)
            else:
                pass

        
    def _render_figure(self, slide: Slide, figure: KFigure):
        #AI模型给的api，都不存在😄
        #from pptx.enum.shapes import MSO_SHAPE,MSO_FLIP
        #需要自己修改xml元素<flip dir='l'></flip>
        #dir='l' 'r' or 't' 'b'
        #picture.element.flip=
        bbox = figure.bbox

        use_pdf=False
        if use_pdf:
            #模拟pdf的渲染方式
            #插入图片，然后可以为上下翻转，左右翻转，旋转一定度数
            #TODO 因为现在的figure使用的是截图的方式
            if figure.flip:
                pass
            elif figure.rotate!=0:
                bbox = self._get_unrotated_bbox(figure)
            else:
                pass
            m = figure.page.to_lt()
            x0,y0,x1,y1 = bbox.transform(m)
            picture = slide.shapes.add_picture(
                str(figure.fullpath),
                Pt(x0), Pt(y0),
                Pt(x1-x0), Pt(y1-y0)
            )
            if figure.flip=='h':
                #deepseek给的api
                #picture.flip(MSO_FLIP.HORIZONTAL)
                #google给的api
                #picture.flip_horizontal=True
                picture.element.flipH=1
            elif figure.flip=='v':
                #picture.flip(MSO_FLIP.VERTICAL)
                #picture.flip_vertical=False
                picture.element.flipV=1
            elif figure.rotate!=0:
                #顺时针旋转多少度
                picture.rotation=figure.rotate
            else:
                pass

        else:
            m = figure.page.to_lt()
            x0,y0,x1,y1 = bbox.transform(m)
            picture = slide.shapes.add_picture(
                str(figure.fullpath),
                Pt(x0), Pt(y0),
                Pt(x1-x0), Pt(y1-y0)
            )

    def _render_table(self, slide: Slide, table: KTable):
        #TODO 也可以按block的方案渲染
        #mode='figure'
        mode='figure'
        if mode=='figure':
            m=table.page.to_lt()
            x0,y0,x1,y1 = table.bbox.transform(m)
            image = table.page.crop(table.bbox)
            assert image is not None
            with io.BytesIO() as fp:
                image.save(fp,format='png')
                fp.seek(0)
                slide.shapes.add_picture(
                    fp,
                    Pt(x0), Pt(y0),
                    Pt(x1-x0), Pt(y1-y0)
                )

        else:
            m=table.page.to_lt()
            x0,y0,x1,y1 = table.bbox.transform(m)
            ptable = slide.shapes.add_table(
                rows=table.row_num, cols=table.col_num,
                left=Pt(x0), top=Pt(y0),
                width=Pt(x1-x0), height=Pt(y1-y0)
            ).table
            
            def get_x_axis()->list[float]:
                axis:list[float] = [None]*table.col_num
                for cell in table.cells:
                    i = cell.col_index
                    # j = cell.col_index+cell.col_span
                    if axis[i] is None:
                        axis[i] = cell.bbox[0]
                axis.append(table.bbox[2])
                return axis

            # Set column widths
            x_axis = get_x_axis()
            for i in range(table.col_num):
                ptable.columns[i].width=Pt(x_axis[i+1]-x_axis[i])
            
            for cell in table.cells:
                pcell = ptable.cell(cell.row_index,cell.col_index)
                #透明
                #pcell.fill.background()
                #对于跨行跨栏的，只需要合并对角就可以
                if cell.col_span>1 or cell.row_span>1:
                    pcell.merge(ptable.cell(cell.row_index+cell.row_span-1,cell.col_index+cell.col_span-1))

                        
                #self._render_block(slide,cell.body,pcell)
                for obj in cell.objects:
                    if isinstance(obj,Text):
                        #文本还是在这里？
                        #remove_all_padding(pcell)
                        self._render_text(slide,obj,text_frame=pcell.text_frame)
                        #TODO 可以根据文本行判断，是否左对齐还是右对齐，还是居中对齐
                        for p in pcell.text_frame.paragraphs:
                            p.alignment=PP_ALIGN.CENTER
                        pcell.text_frame.vertical_anchor=MSO_ANCHOR.MIDDLE
                    elif isinstance(obj,KFigure):
                        self._render_figure(slide,obj)
                    elif isinstance(obj,KFormula):
                        self._render_formula(slide,obj)
                    else:
                        pass
                
    

    def _render_formula(self, slide: Slide, formula: KFormula):
        #可以根据latex渲染为图片，再插入，或者直接使用截图
        #使用渲染的图片可以更加清晰(latex)
        use_latex=False
        if use_latex and formula.latex:
            pass
        else:
            m = formula.page.to_lt()
            x0, y0, x1, y1 = formula.bbox.transform(m)
            # TODO 是否使用相对路径即可？
            slide.shapes.add_picture(
                str(formula.fullpath),
                Pt(x0), Pt(y0),      
                Pt(x1-x0), Pt(y1-y0)
            )
        pass

    def _to_color(self,color:KColor)->RGBColor:
        return RGBColor(color.rgba[0],color.rgba[1],color.rgba[2])
